"""
ROHL Site-Update Builder — dynamic per-space slides + calendar-month rotation.

Model:
  * The deck starts lean (cover + intro slides + feedback). No space slides.
  * Each submitted space gets its OWN slide, created on first submission by
    cloning the matching design prototype (3-column for guest-facing areas,
    2-column for back-of-house/technical). Only filled spaces ever appear.
  * Empty leftover space slides are removed automatically.
  * Slides are kept grouped/ordered by area.
  * Rotation is by CALENDAR MONTH: filling a space again in the same month
    just replaces the current photo; a new month pushes the current photo to
    PREVIOUS and shows the new one as CURRENT. Column headers show real months.

Per-slide state (area / space / current-month) is stored in the slide NOTES,
which never show in the deck.
"""
from pptx import Presentation
from pptx.util import Emu
from pptx.oxml.ns import qn
from PIL import Image
import copy, calendar, re, os

IN = 914400
HERE = os.path.dirname(os.path.abspath(__file__))
TEMPLATE_PATH = os.path.join(HERE, "space_templates.pptx")

# ---- form structure: area -> ordered spaces (drives grouping + order) --------
FORM = {
    "External & Facade": ["Main Entrance / Porch", "Building Exterior", "Parking Area",
                          "Garden / Landscaping", "Driveway", "Boundary Wall", "Signage & Branding"],
    "Public Areas": ["Lobby", "Reception / Front Desk", "Waiting Lounge", "Corridors (Ground Floor)",
                     "Corridors (Upper Floors)", "Elevators", "Staircase", "Public Restrooms", "Luggage Area"],
    "Guest Rooms & Suites": ["Standard Room", "Deluxe Room", "Premium Room", "Junior Suite",
                             "Suite", "Bathroom", "Toilet", "Balcony / Terrace"],
    "F&B Support": ["Restaurant (All Day Dining)", "Bar & Lounge", "Banquet Hall", "Pre-Function Area",
                    "Kitchen", "Room Service Area", "Staff Cafeteria"],
    "Wellness & Recreation": ["Swimming Pool", "Pool Deck", "Gym / Fitness Center", "Spa Reception",
                              "Spa Treatment Rooms", "Steam & Sauna", "Kids Play Area"],
    "Back of House": ["Laundry", "Housekeeping Store", "Staff Locker Room", "General Storage",
                      "Loading Bay", "Linen Room"],
    "Administrative": ["Front Office", "General Manager Office", "HR Room", "Meeting / Training Room", "Finance Office"],
    "Service & Technical": ["Main Electrical Room", "Sub-Electrical Panels", "Plumbing / Water Treatment",
                            "HVAC / Chiller Plant", "Generator Room", "Fire Safety Systems", "CCTV / Security Room"],
}
AREA_LIST = list(FORM.keys())
# guest-facing areas get the 3-column layout (with a design-render column)
THREE_COL_AREAS = {"external & facade", "public areas", "guest rooms & suites",
                   "f&b support", "wellness & recreation"}

def _norm(s):
    return re.sub(r"\s+", " ", (s or "")).strip().lower()

def area_of_space(space):
    ns = _norm(space)
    for area, spaces in FORM.items():
        for sp in spaces:
            if _norm(sp) == ns:
                return area
    return ""

def order_key(area, space):
    a = next((i for i, x in enumerate(AREA_LIST) if _norm(x) == _norm(area)), 99)
    spaces = FORM.get(next((x for x in AREA_LIST if _norm(x) == _norm(area)), ""), [])
    s = next((i for i, x in enumerate(spaces) if _norm(x) == _norm(space)), 99)
    return (a, s)

def layout_for(area, space):
    a = area or area_of_space(space)
    return "3col" if _norm(a) in THREE_COL_AREAS else "2col"

def fmt_month(key):           # "2026-08" -> "AUG 2026"
    try:
        y, m = key.split("-")
        return f"{calendar.month_abbr[int(m)].upper()} {y}"
    except Exception:
        return key or ""

# ---- small shape helpers -----------------------------------------------------
def _txt(sh):
    return sh.text_frame.text.strip() if sh.has_text_frame else ""

def _cx(sh):
    return sh.left + sh.width / 2

def _set_text(sh, txt, hyperlink=None):
    txt = (txt or "").replace("\r\n", "\n").replace("\r", "\n").replace("_x000D_", "")
    p = sh.text_frame.paragraphs[0]
    if p.runs:
        r = p.runs[0]; r.text = txt
        for extra in p.runs[1:]:
            extra.text = ""
    else:
        r = p.add_run(); r.text = txt
    if hyperlink:
        r.hyperlink.address = hyperlink
    return r

def _img_size(path):
    with Image.open(path) as im:
        return im.size

# ---- per-slide state, stored in an OFF-CANVAS marker textbox (never visible) --
def _meta_shape(slide, create=False):
    for sh in slide.shapes:
        if sh.name == "SPACEMETA":
            return sh
    if create:
        tb = slide.shapes.add_textbox(Emu(-2500000), Emu(0), Emu(500000), Emu(300000))
        tb.name = "SPACEMETA"
        return tb
    return None

def get_meta(slide):
    meta = {}
    sh = _meta_shape(slide)
    if sh is not None:
        for line in sh.text_frame.text.splitlines():
            if "=" in line:
                k, v = line.split("=", 1)
                meta[k.strip()] = v.strip()
    return meta

def set_meta(slide, **kw):
    meta = get_meta(slide)
    meta.update({k: v for k, v in kw.items() if v is not None})
    _meta_shape(slide, create=True).text_frame.text = "\n".join(f"{k}={v}" for k, v in meta.items())

# ---- title / subtitle --------------------------------------------------------
def _title_shape(slide):
    for sh in slide.shapes:
        if _txt(sh).startswith("04 ·"):
            return sh
    return None

def _subtitle_shape(slide):
    t = _title_shape(slide)
    if not t:
        return None
    cands = [sh for sh in slide.shapes if sh.has_text_frame and sh is not t
             and sh.top > t.top and abs(sh.left - t.left) < 0.6 * IN and sh.top < 1.6 * IN]
    return min(cands, key=lambda s: s.top) if cands else None

def slide_is_space(slide):
    return _title_shape(slide) is not None

# ---- columns / boxes ---------------------------------------------------------
def find_columns(slide):
    label_cx = {}
    for sh in slide.shapes:
        t = _txt(sh)
        if t.startswith("RENDERED"): label_cx["RENDER"] = _cx(sh)
        elif t.startswith("PREVIOUS"): label_cx["PREV"] = _cx(sh)
        elif t.startswith("CURRENT"): label_cx["CUR"] = _cx(sh)
    rects = [sh for sh in slide.shapes
             if sh.width and sh.height and sh.height > 4 * IN and sh.width > 3 * IN and not _txt(sh)]
    cols = {}
    for key, cx in label_cx.items():
        if rects:
            cols[key] = min(rects, key=lambda r: abs(_cx(r) - cx))
    return cols

def _label_shape(slide, prefix):
    for sh in slide.shapes:
        if _txt(sh).startswith(prefix):
            return sh
    return None

def _placeholder_in(slide, box):
    for sh in slide.shapes:
        if not sh.has_text_frame or sh.height > 1.2 * IN:
            continue
        c = _cx(sh); cy = sh.top + sh.height / 2
        if box.left < c < box.left + box.width and box.top < cy < box.top + box.height:
            if _txt(sh) in ("Photograph awaited", ""):
                return sh
    return None

def _fit(box, iw, ih, inset=0.06 * IN):
    bw, bh = box.width - 2 * inset, box.height - 2 * inset
    s = min(bw / iw, bh / ih)
    nw, nh = iw * s, ih * s
    return int(box.left + inset + (bw - nw) / 2), int(box.top + inset + (bh - nh) / 2), int(nw), int(nh)

def _find_pic(slide, name):
    for sh in slide.shapes:
        if sh.name == name:
            return sh
    return None

def _find_pics(slide, prefix):
    # all pics named e.g. CURIMG, CURIMG_2, CURIMG_3 ...
    return [sh for sh in slide.shapes if sh.name == prefix or sh.name.startswith(prefix + "_")]

def _grid_cells(box, n):
    """Split a box into n tidy cells (1=full, 2=stacked, 3-4=2x2 grid)."""
    from types import SimpleNamespace
    import math
    if n <= 1:
        return [box]
    cols = 1 if n <= 2 else 2
    rows = math.ceil(n / cols)
    gap = int(0.05 * IN)
    cw = (box.width - gap * (cols - 1)) // cols
    ch = (box.height - gap * (rows - 1)) // rows
    cells = []
    for i in range(n):
        r, c = divmod(i, cols)
        cells.append(SimpleNamespace(left=int(box.left + c * (cw + gap)),
                                     top=int(box.top + r * (ch + gap)),
                                     width=int(cw), height=int(ch)))
    return cells

def _place(slide, box, image_path, name):
    iw, ih = _img_size(image_path)
    x, y, w, h = _fit(box, iw, ih)
    pic = slide.shapes.add_picture(image_path, Emu(x), Emu(y), Emu(w), Emu(h))
    pic.name = name
    ph = _placeholder_in(slide, box)
    if ph is not None:
        _set_text(ph, "")
    return pic

def _place_multi(slide, box, paths, prefix):
    """Tile one or more images into a box; first keeps `prefix`, rest get prefix_N."""
    ph = _placeholder_in(slide, box)
    if ph is not None:
        _set_text(ph, "")
    cells = _grid_cells(box, len(paths))
    pics = []
    for i, (cell, p) in enumerate(zip(cells, paths)):
        iw, ih = _img_size(p)
        x, y, w, h = _fit(cell, iw, ih, inset=0.03 * IN)
        pic = slide.shapes.add_picture(p, Emu(x), Emu(y), Emu(w), Emu(h))
        pic.name = prefix if i == 0 else f"{prefix}_{i+1}"
        pics.append(pic)
    return pics

def _has_photo(slide):
    return any(_find_pics(slide, n) for n in ("CURIMG", "PREVIMG", "RENDERIMG"))

# ---- clone a prototype into the hotel deck -----------------------------------
_TEMPLATE = None
def _template():
    global _TEMPLATE
    if _TEMPLATE is None:
        _TEMPLATE = Presentation(TEMPLATE_PATH)
    return _TEMPLATE

def clone_space_slide(prs, layout):
    proto = _template().slides[0 if layout == "3col" else 1]
    new = prs.slides.add_slide(prs.slides[0].slide_layout)
    for sh in list(new.shapes):
        sh._element.getparent().remove(sh._element)
    for sh in proto.shapes:
        new.shapes._spTree.append(copy.deepcopy(sh._element))
    return new

def _replace_text_everywhere(slide, mapping):
    for sh in slide.shapes:
        if not sh.has_text_frame:
            continue
        for para in sh.text_frame.paragraphs:
            for run in para.runs:
                for a, b in mapping.items():
                    if a and a in run.text:
                        run.text = run.text.replace(a, b)

# ---- update actions ----------------------------------------------------------
def set_current_photo(slide, cols, photo_paths, month_key):
    """photo_paths: one path or a list of paths (multiple angles). All are tiled
    into the CURRENT box; on a new month the whole set rotates into PREVIOUS."""
    if isinstance(photo_paths, str):
        photo_paths = [photo_paths]
    photo_paths = [p for p in (photo_paths or []) if p]
    if not photo_paths or "CUR" not in cols:
        return
    stored = get_meta(slide).get("curmonth", "")
    curs = _find_pics(slide, "CURIMG")
    if curs and stored and stored != month_key and "PREV" in cols:
        # new month -> move the whole current set into the PREVIOUS box, retiled
        for pv in _find_pics(slide, "PREVIMG"):
            pv._element.getparent().remove(pv._element)
        cells = _grid_cells(cols["PREV"], len(curs))
        for i, (pic, cell) in enumerate(zip(curs, cells)):
            x, y, w, h = _fit(cell, pic.width, pic.height)
            pic.left, pic.top, pic.width, pic.height = Emu(x), Emu(y), Emu(w), Emu(h)
            pic.name = "PREVIMG" if i == 0 else f"PREVIMG_{i+1}"
        ph = _placeholder_in(slide, cols["PREV"])
        if ph is not None:
            _set_text(ph, "")
        prev_lbl = _label_shape(slide, "PREVIOUS")
        if prev_lbl is not None:
            _set_text(prev_lbl, f"PREVIOUS · {fmt_month(stored)}")
    elif curs and stored == month_key:
        # same month -> replace the current set
        for pic in curs:
            pic._element.getparent().remove(pic._element)
    _place_multi(slide, cols["CUR"], photo_paths, "CURIMG")
    cur_lbl = _label_shape(slide, "CURRENT")
    if cur_lbl is not None:
        _set_text(cur_lbl, f"CURRENT · {fmt_month(month_key)}")
    set_meta(slide, curmonth=month_key)

def set_render(slide, cols, render_path):
    if "RENDER" not in cols or not render_path:
        return
    old = _find_pic(slide, "RENDERIMG")
    if old is not None:
        old._element.getparent().remove(old._element)
    _place(slide, cols["RENDER"], render_path, "RENDERIMG")

def set_videos(slide, videos):
    # single "Site video" link cell (the value box to the right of the label)
    cell = None
    for sh in slide.shapes:
        if sh.has_text_frame and sh.left and abs(sh.left - 2.87 * IN) < 0.4 * IN \
           and sh.top and 8.8 * IN < sh.top < 10.4 * IN:
            cell = sh if cell is None or sh.top < cell.top else cell
    if cell is None:
        return
    v = videos or {}
    url = v.get("walkthrough") or v.get("detail") or v.get("snag") or next((x for x in v.values() if x), None)
    if url:
        _set_text(cell, "Watch ▶", hyperlink=url)

def set_feedback(slide, text):
    lab = _label_shape(slide, "FEEDBACK FOR THIS SPACE")
    if not lab or not text:
        return
    cands = [sh for sh in slide.shapes if sh.has_text_frame and sh is not lab
             and abs(sh.left - lab.left) < 1 * IN and sh.top > lab.top]
    if cands:
        _set_text(max(cands, key=lambda s: s.width * s.height), text)

def set_footer(slide, date_str, completion=""):
    pct = str(completion).strip()
    if pct and not pct.endswith("%"):
        pct = pct + "%"
    pct = pct or "—"
    date_str = date_str or "—"
    for sh in slide.shapes:
        t = _txt(sh)
        if t.startswith("Completion") or t.startswith("Last site visit"):
            _set_text(sh, f"Completion: {pct}   ·   Last site visit: {date_str}   ·   Updated automatically from site form")
            return

# ---- housekeeping ------------------------------------------------------------
def strip_empty_space_slides(prs):
    """Fully delete any space slide that has no photo yet (e.g. the original blanks)."""
    sldIdLst = prs.slides._sldIdLst
    for slide, sid in list(zip(prs.slides, list(sldIdLst))):
        if slide_is_space(slide) and not _has_photo(slide):
            rId = sid.get(qn("r:id"))
            if rId:
                prs.part.drop_rel(rId)
            sldIdLst.remove(sid)


def save_clean(prs, path):
    """Save, then strip orphaned slide/notesSlide parts that python-pptx leaves
    behind after slides are removed (otherwise PowerPoint flags the file for
    repair). Keeps only slides still referenced by presentation.xml.rels."""
    import zipfile, re, os
    prs.save(path)
    z = zipfile.ZipFile(path, "r")
    names = z.namelist()
    presrels = z.read("ppt/_rels/presentation.xml.rels").decode("utf-8")
    keep = set(re.findall(r'Target="(slides/slide\d+\.xml)"', presrels))   # e.g. slides/slide4.xml

    drop = set()
    # unreferenced slide xml + its rels
    for n in names:
        m = re.match(r"ppt/(slides/slide\d+\.xml)$", n)
        if m and m.group(1) not in keep:
            drop.add(n)
            drop.add("ppt/slides/_rels/" + n.split("/")[-1] + ".rels")

    # notesSlides referenced only by dropped slides -> drop them too
    def notes_of(slide_name):
        rp = "ppt/slides/_rels/" + slide_name.split("/")[-1] + ".rels"
        if rp not in names:
            return []
        t = z.read(rp).decode("utf-8")
        return ["ppt/" + x for x in re.findall(r'Target="\.\./(notesSlides/notesSlide\d+\.xml)"', t)]
    kept_notes, dropped_notes = set(), set()
    for n in names:
        if not re.match(r"ppt/slides/slide\d+\.xml$", n):
            continue
        (dropped_notes if n in drop else kept_notes).update(notes_of(n))
    for ns in dropped_notes:
        if ns not in kept_notes:
            drop.add(ns)
            drop.add(ns.replace("notesSlides/", "notesSlides/_rels/") + ".rels")

    # scrub Content_Types overrides for dropped parts
    ct = z.read("[Content_Types].xml").decode("utf-8")
    for d in drop:
        ct = ct.replace(f'<Override PartName="/{d}" ContentType="application/vnd.openxmlformats-officedocument.presentationml.slide+xml"/>', "")
        ct = ct.replace(f'<Override PartName="/{d}" ContentType="application/vnd.openxmlformats-officedocument.presentationml.notesSlide+xml"/>', "")

    tmp = path + ".clean"
    with zipfile.ZipFile(tmp, "w", zipfile.ZIP_DEFLATED) as out:
        for item in z.infolist():
            if item.filename in drop:
                continue
            data = ct.encode("utf-8") if item.filename == "[Content_Types].xml" else z.read(item.filename)
            out.writestr(item, data)
    z.close()
    os.replace(tmp, path)

def reorder_space_slides(prs):
    """Group space slides by area/space order, placed before the feedback slide."""
    sldIdLst = prs.slides._sldIdLst
    pairs = list(zip(prs.slides, list(sldIdLst)))
    space, feedback = [], None
    for slide, sid in pairs:
        t = _title_shape(slide)
        if t is not None:
            m = get_meta(slide)
            space.append((order_key(m.get("area", ""), m.get("space", _txt(t)[5:])), sid))
        elif "Feedback" in " ".join(_txt(s) for s in slide.shapes):
            feedback = sid
    for _, sid in space:
        sldIdLst.remove(sid)
    space.sort(key=lambda x: x[0])
    ids_now = list(sldIdLst)
    at = ids_now.index(feedback) if feedback in ids_now else len(ids_now)
    for off, (_, sid) in enumerate(space):
        sldIdLst.insert(at + off, sid)

def find_or_create(prs, area, space, ctx):
    for slide in prs.slides:
        if slide_is_space(slide) and _norm(get_meta(slide).get("space", "")) == _norm(space):
            return slide
    slide = clone_space_slide(prs, layout_for(area, space))
    _replace_text_everywhere(slide, {
        "ROHL-0002": ctx.get("rohl", ""),
        "Regenta, Jamshedpur": f"{ctx.get('hotel','')}, {ctx.get('city','')}".strip(", "),
        "Jamshedpur": ctx.get("city", ""),
    })
    t = _title_shape(slide)
    if t is not None:
        _set_text(t, f"04 · {space}")
    sub = _subtitle_shape(slide)
    if sub is not None:
        loc = f"{ctx.get('hotel','')}, {ctx.get('city','')}".strip(", ")
        _set_text(sub, f"{area or area_of_space(space)}   ·   {loc}")
    set_meta(slide, area=area or area_of_space(space), space=space)
    return slide

# ---- public entry point ------------------------------------------------------
def apply_submission(prs, sub, context=None):
    ctx = context or {}
    space = sub.get("space", "").strip()
    area = sub.get("area", "").strip() or area_of_space(space)
    month_key = sub.get("month") or ""
    if not space:
        raise ValueError("submission has no space")
    # create/find the slide FIRST (while all slides are present, so a new clone
    # gets a unique part number), fill it, THEN strip the empty leftovers.
    slide = find_or_create(prs, area, space, ctx)
    cols = find_columns(slide)
    if sub.get("render"):
        set_render(slide, cols, sub["render"])
    photos = sub.get("photos") or ([sub["photo"]] if sub.get("photo") else [])
    if photos and "CUR" in cols:
        set_current_photo(slide, cols, photos, month_key)
    if sub.get("videos"):
        set_videos(slide, sub["videos"])
    if sub.get("feedback"):
        set_feedback(slide, sub["feedback"])
    if sub.get("date") or sub.get("completion"):
        set_footer(slide, sub.get("date", ""), sub.get("completion", ""))
    strip_empty_space_slides(prs)
    reorder_space_slides(prs)
    return slide
